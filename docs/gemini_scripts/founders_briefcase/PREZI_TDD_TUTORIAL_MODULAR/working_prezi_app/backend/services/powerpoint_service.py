"""
PowerPoint Service - Core PowerPoint Processing for PrezI
Implements Windows COM automation and cross-platform fallback for PowerPoint file processing
Based on CONSOLIDATED_FOUNDERS_BRIEFCASE.md specifications
"""

import logging
import platform
import tempfile
import hashlib
import uuid
from pathlib import Path
from typing import Dict, List, Any, Optional, Union
from datetime import datetime, timezone

from sqlalchemy.orm import Session
from sqlalchemy import and_

# Windows COM imports (conditionally imported)
if platform.system() == "Windows":
    try:
        import win32com.client
        import pythoncom
        COM_AVAILABLE = True
    except ImportError:
        COM_AVAILABLE = False
        logging.warning("pywin32 not available - PowerPoint COM automation disabled")
else:
    COM_AVAILABLE = False

# Cross-platform fallback imports
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available - cross-platform fallback disabled")

from backend.database.models import FileModel, SlideModel, ElementModel
from backend.core.config import get_settings

logger = logging.getLogger(__name__)

class PowerPointProcessingError(Exception):
    """Custom exception for PowerPoint processing errors"""
    pass

class PowerPointProcessor:
    """
    PowerPoint file processor with COM automation and cross-platform fallback
    Implements the core PowerPoint integration from CONSOLIDATED_FOUNDERS_BRIEFCASE.md
    """
    
    def __init__(self):
        self.settings = get_settings()
        self.use_com = COM_AVAILABLE and platform.system() == "Windows" and self.settings.supports_com
        
        if not self.use_com and not PPTX_AVAILABLE:
            raise PowerPointProcessingError(
                "Neither COM automation nor python-pptx is available. "
                "Please install python-pptx for cross-platform support or run on Windows with PowerPoint installed."
            )
        
        self.processor_type = "COM" if self.use_com else "python-pptx"
        logger.info(f"PowerPoint processor initialized using: {self.processor_type}")
    
    def process_file(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Process PowerPoint file and extract slide content
        
        Args:
            file_path: Path to PowerPoint file
            
        Returns:
            Dictionary containing slide data and metadata
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        
        logger.info(f"Processing PowerPoint file: {file_path.name}")
        
        try:
            if self.use_com:
                logger.info("Using Windows COM automation for PowerPoint processing")
                return self._process_with_com(file_path)
            else:
                logger.info("Using python-pptx for PowerPoint processing")
                return self._process_with_pptx(file_path)
                
        except Exception as e:
            logger.error(f"Failed to process PowerPoint file {file_path}: {e}")
            raise PowerPointProcessingError(f"PowerPoint processing failed: {str(e)}") from e
    
    def _process_with_com(self, file_path: Path) -> Dict[str, Any]:
        """Process using Windows COM automation - implements CONSOLIDATED_FOUNDERS_BRIEFCASE.md specs"""
        
        # Initialize COM
        pythoncom.CoInitialize()
        powerpoint = None
        presentation = None
        
        try:
            # Create PowerPoint application
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = False  # Keep hidden for processing
            
            # Open presentation
            presentation = powerpoint.Presentations.Open(
                str(file_path.absolute()),
                ReadOnly=True,
                Untitled=True,
                WithWindow=False
            )
            
            slide_count = presentation.Slides.Count
            logger.info(f"Found {slide_count} slides in presentation")
            
            slides_data = []
            
            # Process each slide
            for slide_num in range(1, slide_count + 1):
                slide = presentation.Slides.Item(slide_num)
                slide_data = self._extract_slide_data_com(slide, slide_num, file_path)
                slides_data.append(slide_data)
                logger.debug(f"Processed slide {slide_num}/{slide_count}")
            
            # Extract presentation metadata
            presentation_info = {
                "slide_count": slide_count,
                "slides": slides_data,
                "processor": "COM",
                "presentation_name": presentation.Name,
                "has_macros": presentation.HasVBProject if hasattr(presentation, 'HasVBProject') else False,
                "creation_date": self._get_creation_date_com(presentation)
            }
            
            logger.info(f"Successfully processed {slide_count} slides using COM")
            return presentation_info
            
        finally:
            # Clean up resources
            if presentation:
                try:
                    presentation.Close()
                except:
                    pass
            
            if powerpoint:
                try:
                    powerpoint.Quit()
                except:
                    pass
            
            pythoncom.CoUninitialize()
    
    def _extract_slide_data_com(self, slide, slide_number: int, file_path: Path) -> Dict[str, Any]:
        """Extract comprehensive data from a single slide using COM"""
        
        # Initialize slide data
        title_text = ""
        content_text = ""
        elements = []
        
        # Process all shapes on the slide
        try:
            shapes_count = slide.Shapes.Count
            for shape_index in range(1, shapes_count + 1):
                try:
                    shape = slide.Shapes.Item(shape_index)
                    element_data = self._extract_shape_data_com(shape, shape_index)
                    
                    if element_data:
                        # Determine if this is title or content based on position and order
                        if shape_index == 1 and element_data["type"] == "text" and not title_text:
                            title_text = element_data["content"]
                        elif element_data["type"] == "text":
                            content_text += f"{element_data['content']}\n"
                        
                        elements.append(element_data)
                
                except Exception as e:
                    logger.warning(f"Failed to extract shape {shape_index} from slide {slide_number}: {e}")
                    continue
        
        except Exception as e:
            logger.warning(f"Failed to process shapes for slide {slide_number}: {e}")
        
        # Extract speaker notes
        notes_text = self._extract_notes_com(slide, slide_number)
        
        # Generate thumbnail and full image
        thumbnail_path = self._generate_thumbnail_com(slide, slide_number, file_path)
        full_image_path = self._generate_full_image_com(slide, slide_number, file_path)
        
        return {
            "slide_number": slide_number,
            "title": title_text.strip(),
            "content": content_text.strip(),
            "notes": notes_text.strip(),
            "thumbnail_path": thumbnail_path,
            "full_image_path": full_image_path,
            "elements": elements,
            "layout_name": self._get_layout_name_com(slide),
            "background_type": self._get_background_type_com(slide)
        }
    
    def _extract_shape_data_com(self, shape, shape_index: int) -> Optional[Dict[str, Any]]:
        """Extract data from a single shape using COM"""
        
        try:
            element_data = {
                "index": shape_index,
                "type": self._determine_element_type_com(shape),
                "content": "",
                "position": {
                    "x": float(shape.Left) if hasattr(shape, 'Left') else 0,
                    "y": float(shape.Top) if hasattr(shape, 'Top') else 0,
                    "width": float(shape.Width) if hasattr(shape, 'Width') else 0,
                    "height": float(shape.Height) if hasattr(shape, 'Height') else 0,
                    "z_order": shape.ZOrderPosition if hasattr(shape, 'ZOrderPosition') else shape_index
                }
            }
            
            # Extract text content if shape has text
            if hasattr(shape, 'HasTextFrame') and shape.HasTextFrame:
                try:
                    text_content = shape.TextFrame.TextRange.Text.strip()
                    element_data["content"] = text_content
                except:
                    pass
            
            # Extract additional properties based on shape type
            if element_data["type"] == "chart":
                element_data["chart_type"] = self._get_chart_type_com(shape)
            elif element_data["type"] == "image":
                element_data["image_name"] = getattr(shape, 'Name', '')
            elif element_data["type"] == "table":
                element_data["table_info"] = self._get_table_info_com(shape)
            
            return element_data if element_data["content"] or element_data["type"] != "text" else None
            
        except Exception as e:
            logger.warning(f"Failed to extract shape data: {e}")
            return None
    
    def _extract_notes_com(self, slide, slide_number: int) -> str:
        """Extract speaker notes from slide using COM"""
        notes_text = ""
        
        try:
            notes_slide = slide.NotesPage
            shapes_count = notes_slide.Shapes.Count
            
            for shape_index in range(1, shapes_count + 1):
                try:
                    shape = notes_slide.Shapes.Item(shape_index)
                    
                    if hasattr(shape, 'HasTextFrame') and shape.HasTextFrame:
                        notes_content = shape.TextFrame.TextRange.Text.strip()
                        if notes_content:
                            notes_text += f"{notes_content}\n"
                except:
                    continue
                        
        except Exception as e:
            logger.warning(f"Failed to extract notes for slide {slide_number}: {e}")
        
        return notes_text.strip()
    
    def _generate_thumbnail_com(self, slide, slide_number: int, file_path: Path) -> str:
        """Generate thumbnail image using COM"""
        try:
            # Create thumbnail directory
            thumbnail_dir = Path(self.settings.thumbnail_dir)
            thumbnail_dir.mkdir(parents=True, exist_ok=True)
            
            # Generate unique filename based on file and slide
            file_hash = hashlib.md5(str(file_path).encode()).hexdigest()[:8]
            timestamp = int(datetime.now().timestamp())
            thumbnail_name = f"{file_hash}_slide_{slide_number}_{timestamp}.png"
            thumbnail_path = thumbnail_dir / thumbnail_name
            
            # Export slide as thumbnail
            slide.Export(
                str(thumbnail_path),
                "PNG",
                ScaleWidth=300,  # Standard thumbnail width
                ScaleHeight=225  # Standard thumbnail height (4:3 ratio)
            )
            
            if thumbnail_path.exists():
                logger.debug(f"Generated thumbnail: {thumbnail_path}")
                return str(thumbnail_path)
            else:
                logger.warning(f"Thumbnail generation failed for slide {slide_number}")
                return ""
            
        except Exception as e:
            logger.error(f"Failed to generate thumbnail for slide {slide_number}: {e}")
            return ""
    
    def _generate_full_image_com(self, slide, slide_number: int, file_path: Path) -> str:
        """Generate full-size image using COM"""
        try:
            # Create image directory
            image_dir = Path(self.settings.upload_dir).parent / "images"
            image_dir.mkdir(parents=True, exist_ok=True)
            
            # Generate unique filename
            file_hash = hashlib.md5(str(file_path).encode()).hexdigest()[:8]
            timestamp = int(datetime.now().timestamp())
            image_name = f"{file_hash}_slide_{slide_number}_{timestamp}_full.png"
            image_path = image_dir / image_name
            
            # Export slide as high-resolution image
            slide.Export(
                str(image_path),
                "PNG",
                ScaleWidth=1920,  # High resolution width
                ScaleHeight=1440  # High resolution height
            )
            
            if image_path.exists():
                logger.debug(f"Generated full image: {image_path}")
                return str(image_path)
            else:
                logger.warning(f"Full image generation failed for slide {slide_number}")
                return ""
            
        except Exception as e:
            logger.error(f"Failed to generate full image for slide {slide_number}: {e}")
            return ""
    
    def _determine_element_type_com(self, shape) -> str:
        """Determine element type from COM shape"""
        try:
            # PowerPoint shape type constants
            shape_type = shape.Type
            
            if shape_type == 1:  # msoAutoShape or msoTextBox
                if hasattr(shape, 'HasTextFrame') and shape.HasTextFrame:
                    return "text"
                else:
                    return "shape"
            elif shape_type == 13:  # msoPicture
                return "image"
            elif shape_type == 3:  # msoChart
                return "chart"
            elif shape_type == 19:  # msoTable
                return "table"
            elif shape_type == 7:  # msoMedia (video/audio)
                return "video"
            elif shape_type == 16:  # msoSmartArt
                return "smartart"
            else:
                return "other"
                
        except Exception:
            return "other"
    
    def _get_layout_name_com(self, slide) -> str:
        """Get slide layout name using COM"""
        try:
            return slide.Layout.Name
        except:
            return "Unknown Layout"
    
    def _get_background_type_com(self, slide) -> str:
        """Get slide background type using COM"""
        try:
            if hasattr(slide.Background, 'Type'):
                bg_type = slide.Background.Type
                if bg_type == 1:
                    return "solid"
                elif bg_type == 2:
                    return "gradient"
                elif bg_type == 3:
                    return "texture"
                elif bg_type == 4:
                    return "pattern"
                elif bg_type == 5:
                    return "picture"
            return "default"
        except:
            return "unknown"
    
    def _get_chart_type_com(self, shape) -> str:
        """Get chart type for chart shapes"""
        try:
            if hasattr(shape, 'Chart'):
                chart_type = shape.Chart.ChartType
                # Map common chart types
                chart_types = {
                    51: "column",
                    52: "clustered_column",
                    4: "line",
                    5: "pie",
                    15: "area",
                    65: "bar"
                }
                return chart_types.get(chart_type, "unknown")
        except:
            pass
        return "unknown"
    
    def _get_table_info_com(self, shape) -> Dict[str, Any]:
        """Get table information for table shapes"""
        try:
            if hasattr(shape, 'Table'):
                table = shape.Table
                return {
                    "rows": table.Rows.Count,
                    "columns": table.Columns.Count
                }
        except:
            pass
        return {"rows": 0, "columns": 0}
    
    def _get_creation_date_com(self, presentation) -> Optional[str]:
        """Get presentation creation date using COM"""
        try:
            if hasattr(presentation, 'BuiltInDocumentProperties'):
                creation_date = presentation.BuiltInDocumentProperties.Item("Creation Date").Value
                if creation_date:
                    return str(creation_date)
        except:
            pass
        return None
    
    def _process_with_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Process using python-pptx library (cross-platform fallback)"""
        
        try:
            presentation = Presentation(str(file_path))
            slide_count = len(presentation.slides)
            
            logger.info(f"Found {slide_count} slides in presentation (python-pptx)")
            
            slides_data = []
            
            # Process each slide
            for slide_index, slide in enumerate(presentation.slides):
                slide_data = self._extract_slide_data_pptx(slide, slide_index + 1, file_path)
                slides_data.append(slide_data)
                logger.debug(f"Processed slide {slide_index + 1}/{slide_count}")
            
            logger.info(f"Successfully processed {slide_count} slides using python-pptx")
            
            return {
                "slide_count": slide_count,
                "slides": slides_data,
                "processor": "python-pptx",
                "presentation_name": file_path.stem,
                "has_macros": False,  # python-pptx can't detect macros
                "creation_date": None
            }
            
        except Exception as e:
            raise PowerPointProcessingError(f"python-pptx processing failed: {str(e)}") from e
    
    def _extract_slide_data_pptx(self, slide, slide_number: int, file_path: Path) -> Dict[str, Any]:
        """Extract data from a single slide using python-pptx"""
        
        title_text = ""
        content_text = ""
        elements = []
        
        # Extract text from shapes
        for shape_index, shape in enumerate(slide.shapes):
            try:
                element_data = self._extract_shape_data_pptx(shape, shape_index)
                
                if element_data:
                    # First text shape is usually title
                    if not title_text and element_data["type"] == "text":
                        title_text = element_data["content"]
                    elif element_data["type"] == "text":
                        content_text += f"{element_data['content']}\n"
                    
                    elements.append(element_data)
                    
            except Exception as e:
                logger.warning(f"Failed to extract shape {shape_index} from slide {slide_number}: {e}")
                continue
        
        # Extract speaker notes
        notes_text = ""
        try:
            if hasattr(slide, 'notes_slide') and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
        except Exception as e:
            logger.warning(f"Failed to extract notes for slide {slide_number}: {e}")
        
        # Note: python-pptx doesn't support image export
        # We'll use placeholder paths for now
        return {
            "slide_number": slide_number,
            "title": title_text,
            "content": content_text.strip(),
            "notes": notes_text,
            "thumbnail_path": "",  # Would need additional library for image export
            "full_image_path": "",
            "elements": elements,
            "layout_name": slide.slide_layout.name if hasattr(slide, 'slide_layout') else "Unknown",
            "background_type": "unknown"
        }
    
    def _extract_shape_data_pptx(self, shape, shape_index: int) -> Optional[Dict[str, Any]]:
        """Extract data from a single shape using python-pptx"""
        
        try:
            element_data = {
                "index": shape_index,
                "type": self._determine_element_type_pptx(shape),
                "content": "",
                "position": {
                    "x": float(shape.left) if hasattr(shape, 'left') else 0,
                    "y": float(shape.top) if hasattr(shape, 'top') else 0,
                    "width": float(shape.width) if hasattr(shape, 'width') else 0,
                    "height": float(shape.height) if hasattr(shape, 'height') else 0,
                    "z_order": 0  # python-pptx doesn't provide z-order
                }
            }
            
            # Extract text content
            if hasattr(shape, "text") and shape.text.strip():
                element_data["content"] = shape.text.strip()
            
            return element_data if element_data["content"] or element_data["type"] != "text" else None
            
        except Exception as e:
            logger.warning(f"Failed to extract shape data: {e}")
            return None
    
    def _determine_element_type_pptx(self, shape) -> str:
        """Determine element type from python-pptx shape"""
        try:
            if hasattr(shape, 'shape_type'):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    return "image"
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    return "chart"
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    return "table"
                elif hasattr(shape, 'text'):
                    return "text"
                else:
                    return "shape"
            else:
                return "other"
                
        except Exception:
            return "other"

class PowerPointService:
    """
    PowerPoint service for file processing and database integration
    Implements the PowerPoint processing pipeline from CONSOLIDATED_FOUNDERS_BRIEFCASE.md
    """
    
    def __init__(self, db_session: Session):
        self.db = db_session
        try:
            self.processor = PowerPointProcessor()
        except PowerPointProcessingError as e:
            logger.error(f"Failed to initialize PowerPoint processor: {e}")
            raise
    
    def process_file(self, file_id: str) -> Dict[str, Any]:
        """
        Process a PowerPoint file and save results to database
        
        Args:
            file_id: Database ID of file to process
            
        Returns:
            Processing result dictionary
        """
        # Get file record
        file_record = self.db.query(FileModel).filter(FileModel.id == file_id).first()
        if not file_record:
            raise PowerPointProcessingError(f"File {file_id} not found in database")
        
        # Update processing status
        file_record.processed = False
        file_record.processing_error = None
        self.db.commit()
        
        try:
            logger.info(f"Starting PowerPoint processing for file: {file_record.filename}")
            
            # Process file using PowerPoint processor
            result = self.processor.process_file(file_record.file_path)
            
            # Save slides to database
            self._save_slides_to_database(file_record, result["slides"])
            
            # Update file record with success
            file_record.processed = True
            file_record.slide_count = result["slide_count"]
            file_record.processing_error = None
            
            # Update AI analysis if available
            if "has_macros" in result:
                file_record.ai_analysis = {
                    "processor": result["processor"],
                    "has_macros": result["has_macros"],
                    "creation_date": result.get("creation_date"),
                    "processing_timestamp": datetime.now().isoformat()
                }
            
            self.db.commit()
            
            logger.info(f"Successfully processed {result['slide_count']} slides from {file_record.filename}")
            
            return {
                "success": True,
                "file_id": file_id,
                "filename": file_record.filename,
                "slide_count": result["slide_count"],
                "processor": result["processor"]
            }
            
        except Exception as e:
            # Update file record with error
            file_record.processed = False
            file_record.processing_error = str(e)
            self.db.commit()
            
            logger.error(f"Failed to process PowerPoint file {file_record.filename}: {e}")
            
            return {
                "success": False,
                "file_id": file_id,
                "filename": file_record.filename,
                "error": str(e)
            }
    
    def bulk_process_files(self, file_ids: List[str]) -> Dict[str, Any]:
        """
        Process multiple PowerPoint files in batch
        
        Args:
            file_ids: List of file IDs to process
            
        Returns:
            Bulk processing results
        """
        results = {
            "successful": [],
            "failed": [],
            "total_files": len(file_ids),
            "total_slides": 0
        }
        
        logger.info(f"Starting bulk processing of {len(file_ids)} PowerPoint files")
        
        for file_id in file_ids:
            try:
                result = self.process_file(file_id)
                
                if result["success"]:
                    results["successful"].append(result)
                    results["total_slides"] += result["slide_count"]
                else:
                    results["failed"].append(result)
                    
            except Exception as e:
                logger.error(f"Unexpected error processing file {file_id}: {e}")
                results["failed"].append({
                    "success": False,
                    "file_id": file_id,
                    "error": str(e)
                })
        
        logger.info(f"Bulk processing complete: {len(results['successful'])} successful, {len(results['failed'])} failed")
        
        return results
    
    def _save_slides_to_database(self, file_record: FileModel, slides_data: List[Dict[str, Any]]):
        """Save extracted slide data to database"""
        
        # Clear existing slides for this file
        self.db.query(SlideModel).filter(SlideModel.file_id == file_record.id).delete()
        
        for slide_data in slides_data:
            # Create slide record
            slide = SlideModel(
                id=str(uuid.uuid4()),
                file_id=file_record.id,
                slide_number=slide_data["slide_number"],
                title=slide_data["title"][:500] if slide_data["title"] else None,  # Limit title length
                notes=slide_data["notes"],
                slide_type=self._classify_slide_type(slide_data),
                thumbnail_path=slide_data["thumbnail_path"],
                full_image_path=slide_data["full_image_path"],
                ai_analysis={
                    "layout_name": slide_data["layout_name"],
                    "background_type": slide_data["background_type"],
                    "element_count": len(slide_data["elements"]),
                    "extracted_timestamp": datetime.now().isoformat()
                }
            )
            
            self.db.add(slide)
            self.db.flush()  # Get slide ID
            
            # Create element records
            for element_data in slide_data["elements"]:
                element = ElementModel(
                    id=str(uuid.uuid4()),
                    slide_id=slide.id,
                    element_type=element_data["type"],
                    content=element_data["content"][:1000] if element_data["content"] else None,  # Limit content
                    position_x=element_data["position"]["x"],
                    position_y=element_data["position"]["y"],
                    width=element_data["position"]["width"],
                    height=element_data["position"]["height"],
                    ai_analysis={
                        "index": element_data["index"],
                        "z_order": element_data["position"].get("z_order", 0),
                        **{k: v for k, v in element_data.items() if k not in ["index", "type", "content", "position"]}
                    }
                )
                
                self.db.add(element)
        
        self.db.commit()
    
    def _classify_slide_type(self, slide_data: Dict[str, Any]) -> str:
        """Classify slide type based on content and structure"""
        
        # Simple heuristic-based classification
        title = slide_data.get("title", "").lower()
        content = slide_data.get("content", "").lower()
        element_count = len(slide_data.get("elements", []))
        elements = slide_data.get("elements", [])
        
        # Check for specific slide types
        if any(keyword in title for keyword in ["introduction", "welcome", "agenda", "overview"]):
            return "title"
        elif any(keyword in title for keyword in ["conclusion", "summary", "thank you", "questions"]):
            return "conclusion"
        elif any(element["type"] == "chart" for element in elements):
            return "chart"
        elif any(element["type"] == "image" for element in elements):
            return "image"
        elif any(element["type"] == "table" for element in elements):
            return "table"
        elif element_count > 3:
            return "content"
        else:
            return "unknown"
    
    def get_processing_status(self, file_id: str) -> Dict[str, Any]:
        """Get processing status for a file"""
        file_record = self.db.query(FileModel).filter(FileModel.id == file_id).first()
        if not file_record:
            raise PowerPointProcessingError(f"File {file_id} not found")
        
        return {
            "file_id": file_id,
            "filename": file_record.filename,
            "processed": file_record.processed,
            "slide_count": file_record.slide_count,
            "processing_error": file_record.processing_error,
            "processor_available": self.processor.processor_type
        }
    
    def reprocess_file(self, file_id: str, force: bool = False) -> Dict[str, Any]:
        """Reprocess a PowerPoint file (useful for failed processing)"""
        file_record = self.db.query(FileModel).filter(FileModel.id == file_id).first()
        if not file_record:
            raise PowerPointProcessingError(f"File {file_id} not found")
        
        if file_record.processed and not force:
            return {
                "success": False,
                "error": "File already processed. Use force=True to reprocess."
            }
        
        return self.process_file(file_id)

# Example usage and testing
if __name__ == "__main__":
    # This would be used for testing the PowerPoint service
    import tempfile
    
    # Create a mock database session for testing
    class MockSession:
        def query(self, model):
            return self
        def filter(self, condition):
            return self
        def first(self):
            return None
        def add(self, obj):
            pass
        def commit(self):
            pass
        def flush(self):
            pass
        def delete(self):
            pass
    
    # Test processor initialization
    try:
        processor = PowerPointProcessor()
        print(f"✅ PowerPoint processor initialized: {processor.processor_type}")
        
        service = PowerPointService(MockSession())
        print("✅ PowerPoint service initialized successfully")
        
    except PowerPointProcessingError as e:
        print(f"❌ Failed to initialize PowerPoint service: {e}")
    except Exception as e:
        print(f"❌ Unexpected error: {e}")