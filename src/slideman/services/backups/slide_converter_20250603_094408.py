# src/slideman/services/slide_converter.py

import logging
import time # For potential delays/timeouts
import traceback
from pathlib import Path
from typing import Optional

# --- Qt Imports ---
from PySide6.QtCore import QObject, QRunnable, Signal, Slot, QSize, Qt
from PySide6.QtGui import QPixmap # For thumbnail generation via Qt

# --- External Libraries ---
# Import COM libraries carefully
try:
    import pythoncom
    import win32com.client
    import pywintypes
    HAS_COM = True
except ImportError:
    HAS_COM = False
    logging.error("pywin32 library not found. PowerPoint COM automation will not work.")
    # Define dummy types for type hinting if needed
    pywintypes = type('pywintypes', (object,), {'com_error': Exception})()


# Import python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Emu # Use Emu for dimensions
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logging.error("python-pptx library not found. Shape extraction will not work.")
    # Define dummy types
    MSO_SHAPE_TYPE = type('MSO_SHAPE_TYPE', (object,), {'PLACEHOLDER': 0, 'PICTURE': 1, 'CHART': 2, 'TABLE': 3, 'TEXT_BOX': 4, 'GROUP': 5})()


# Import Pillow (Optional - can use QPixmap for resizing too)
# try:
#     from PIL import Image
#     HAS_PIL = True
# except ImportError:
#     HAS_PIL = False
#     logging.warning("Pillow library not found. Thumbnail generation might fallback to Qt.")

# --- Local Imports ---
from .database_worker import DatabaseWorker  # Thread-safe DB access
from .exceptions import (
    DatabaseError, PowerPointError, COMInitializationError,
    PresentationAccessError, SlideExportError
)

# --- Constants ---
THUMBNAIL_HEIGHT = 200 # Target height in pixels for thumbnails

logger = logging.getLogger(__name__)

class SlideConverterSignals(QObject):
    """Defines signals emitted by the SlideConverter worker."""
    progress = Signal(int, int, int) # file_id, current_slide_index (1-based), total_slides
    finished = Signal(int)           # file_id
    error = Signal(int, str)         # file_id, error_message

class SlideConverter(QRunnable):
    """
    Worker thread to convert a PowerPoint file's slides to images,
    extract shape data, and store results in the database.
    """
    def __init__(self, file_id: int, file_path: Path, db_path: Path, signals: SlideConverterSignals):
        super().__init__()
        if not HAS_COM or not HAS_PPTX:
             raise ImportError("Missing required libraries (pywin32 or python-pptx) for SlideConverter.")

        self.signals = signals
        self.file_id = file_id
        self.file_path = file_path
        # Store the database path for worker thread connection
        self.db_path = db_path
        self.is_cancelled = False

        # Determine base path for converted data within project structure
        # Assumes file_path is INSIDE the project structure already
        self.project_root = file_path.parent.parent # Go up two levels: file -> project_folder
        self.converted_data_dir = self.project_root / "converted_data" / str(self.file_id)


    def map_shape_type(self, pptx_shape_type) -> str:
         """Maps python-pptx MSO_SHAPE_TYPE to simple strings."""
         # This mapping might need refinement based on testing different shapes
         if pptx_shape_type in (MSO_SHAPE_TYPE.PLACEHOLDER, MSO_SHAPE_TYPE.AUTO_SHAPE,
                                MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.GROUP): # Treat Group shapes? maybe skip?
             return "SHAPE"
         elif pptx_shape_type == MSO_SHAPE_TYPE.PICTURE:
             return "PICTURE"
         elif pptx_shape_type == MSO_SHAPE_TYPE.CHART:
             return "CHART"
         elif pptx_shape_type == MSO_SHAPE_TYPE.TABLE:
             return "TABLE"
         elif pptx_shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
              # Check if it actually contains text, otherwise might be just a shape
              # shape_obj = # Need the actual shape object here
              # if shape_obj.has_text_frame and shape_obj.text.strip():
              #      return "TEXT"
              # else:
              #      return "SHAPE"
              return "TEXT" # Simplify for now
         # Add other types if needed (MEDIA, etc.)
         else:
             return "UNKNOWN"

    @Slot()
    def run(self):
        """Main execution logic running in the background thread."""
        logger.info(f"[Converter FID:{self.file_id}] Starting conversion for: {self.file_path}")

        # Create a thread-safe database worker
        db_worker = DatabaseWorker(self.db_path)
        
        ppt_app = None
        presentation_com = None
        total_slides = 0
        processed_slides = 0
        encountered_error = False
        error_message = ""
        com_initialized = False

        try:
            # --- 1. Initialize COM ---
            try:
                pythoncom.CoInitialize()
                com_initialized = True
                logger.debug(f"[Converter FID:{self.file_id}] COM Initialized.")
            except Exception as e:
                raise COMInitializationError(f"Failed to initialize COM: {e}") from e

            # --- 2. Update conversion status ---
            try:
                db_worker.update_file_conversion_status(self.file_id, "In Progress")
            except DatabaseError as e:
                logger.error(f"[Converter FID:{self.file_id}] Failed to update conversion status: {e}")
                # Continue anyway, non-critical

        try:
            # --- 2. Create Output Directories ---
            # Example: ProjectRoot/converted_data/FILE_ID/
            self.converted_data_dir.mkdir(parents=True, exist_ok=True)
            logger.debug(f"[Converter FID:{self.file_id}] Ensured output directory exists: {self.converted_data_dir}")

            # --- 3. Open Presentations (COM & python-pptx) ---
            logger.debug(f"[Converter FID:{self.file_id}] Opening presentation via COM...")
            try:
                # Use Dispatch for robustness against PowerPoint already running?
                # Consider using EnsureDispatch if you want to reuse an existing instance (more complex)
                ppt_app = win32com.client.Dispatch("PowerPoint.Application")
                # Optional: Make invisible? ppt_app.Visible = False might not work reliably with Dispatch
                presentation_com = ppt_app.Presentations.Open(str(self.file_path), ReadOnly=True, WithWindow=False)
                total_slides = len(presentation_com.Slides)
                logger.info(f"[Converter FID:{self.file_id}] Opened via COM. Found {total_slides} slides.")
            except pywintypes.com_error as e:
                 logger.error(f"[Converter FID:{self.file_id}] COM Error opening presentation: {e}", exc_info=True)
                 raise RuntimeError(f"PowerPoint COM Error: {e}") # Re-raise to be caught below

            logger.debug(f"[Converter FID:{self.file_id}] Opening presentation via python-pptx...")
            try:
                 presentation_pptx = Presentation(self.file_path)
                 if total_slides != len(presentation_pptx.slides):
                      logger.warning(f"[Converter FID:{self.file_id}] Slide count mismatch! COM: {total_slides}, python-pptx: {len(presentation_pptx.slides)}. Using COM count.")
                 # Use total_slides from COM as primary count
            except Exception as e:
                 logger.error(f"[Converter FID:{self.file_id}] Error opening with python-pptx: {e}", exc_info=True)
                 raise RuntimeError(f"python-pptx Error: {e}")

            if total_slides == 0:
                 logger.warning(f"[Converter FID:{self.file_id}] Presentation has no slides. Marking as complete.")
                 thread_local_db.update_file_conversion_status(self.file_id, 'Completed')
                 self.signals.finished.emit(self.file_id)
                 return # Nothing more to do


            # --- 4. Iterate Through Slides ---
            for i in range(total_slides):
                if self.is_cancelled:
                    logger.info(f"[Converter FID:{self.file_id}] Conversion cancelled.")
                    encountered_error = True
                    error_message = "Cancelled"
                    break # Exit loop

                slide_index_1based = i + 1
                logger.debug(f"[Converter FID:{self.file_id}] Processing slide {slide_index_1based}/{total_slides}")
                slide_processed_successfully = False

                try:
                    slide_com = presentation_com.Slides(slide_index_1based) # Get COM slide object (1-based index)
                    # Ensure we have a matching python-pptx slide object if counts match
                    slide_pptx = None
                    if i < len(presentation_pptx.slides):
                        slide_pptx = presentation_pptx.slides[i] # 0-based index
                    else:
                        logger.warning(f"[Converter FID:{self.file_id}] Skipping shape extraction for slide {slide_index_1based} due to count mismatch.")


                    # --- 4a. Export Full-Res Image (COM) ---
                    img_filename = f"image_{slide_index_1based}.png"
                    img_full_path = self.converted_data_dir / img_filename
                    logger.debug(f"[Converter FID:{self.file_id}] Exporting full-res image to {img_full_path}")
                    # TODO: What width/height for export? Use slide dimensions? Default?
                    # Let's use default for now. Can add Width/Height args if needed.
                    slide_com.Export(str(img_full_path), "PNG")
                    image_rel_path = img_full_path.relative_to(self.project_root).as_posix() # Store POSIX path

                    # --- 4b. Create Thumbnail (Using QPixmap) ---
                    thumb_filename = f"thumb_{slide_index_1based}.png"
                    thumb_full_path = self.converted_data_dir / thumb_filename
                    logger.debug(f"[Converter FID:{self.file_id}] Creating thumbnail image at {thumb_full_path}")
                    pixmap = QPixmap(str(img_full_path))
                    if pixmap.isNull():
                         logger.error(f"[Converter FID:{self.file_id}] Failed to load exported image for thumbnailing: {img_full_path}")
                         raise RuntimeError("Failed to load exported image.")

                    # Scale pixmap to desired height, keeping aspect ratio
                    thumb_pixmap = pixmap.scaledToHeight(THUMBNAIL_HEIGHT, Qt.TransformationMode.SmoothTransformation)
                    if not thumb_pixmap.save(str(thumb_full_path), "PNG"):
                         logger.error(f"[Converter FID:{self.file_id}] Failed to save thumbnail image: {thumb_full_path}")
                         # Continue without thumbnail? Or fail? Let's continue but log error.
                         thumb_rel_path = None
                    else:
                         thumb_rel_path = thumb_full_path.relative_to(self.project_root).as_posix()

                    # --- 4c. Add/Update Slide Record in DB ---
                    logger.debug(f"[Converter FID:{self.file_id}] Updating database for slide {slide_index_1based}")
                    slide_id = thread_local_db.add_slide(self.file_id, slide_index_1based, thumb_rel_path, image_rel_path)
                    if slide_id is None:
                        raise RuntimeError("Failed to add/update slide record in database.")

                    # --- 4d. Extract and Store Elements (python-pptx) ---
                    if slide_pptx:
                        logger.debug(f"[Converter FID:{self.file_id}] Extracting elements for slide {slide_index_1based}")
                        # Delete existing elements for this slide before adding new ones
                        if not thread_local_db.delete_elements_for_slide(slide_id):
                            logger.warning(f"[Converter FID:{self.file_id}] Failed to delete existing elements for SlideID {slide_id}. New elements might be duplicates.")

                        element_count = 0
                        for shape in slide_pptx.shapes:
                            # Basic check if shape has usable geometry
                            if not hasattr(shape, 'shape_type') or not hasattr(shape, 'left') or \
                               shape.left is None or shape.top is None or \
                               shape.width is None or shape.height is None:
                                 logger.debug(f"[Converter FID:{self.file_id}] Skipping shape without type/geometry: {getattr(shape, 'name', 'Unnamed')}")
                                 continue

                            element_type_str = self.map_shape_type(shape.shape_type)
                            bbox_x = float(shape.left.emu)
                            bbox_y = float(shape.top.emu)
                            bbox_w = float(shape.width.emu)
                            bbox_h = float(shape.height.emu)
                            

                            # Add element to DB
                            element_id = thread_local_db.add_element(slide_id, element_type_str, bbox_x, bbox_y, bbox_w, bbox_h)
                            if element_id:
                                 element_count += 1
                            else:
                                 logger.warning(f"[Converter FID:{self.file_id}] Failed to add element DB record for shape: {getattr(shape, 'name', 'Unnamed')}")
                        logger.debug(f"[Converter FID:{self.file_id}] Added {element_count} elements for SlideID {slide_id}")
                    else:
                         logger.warning(f"[Converter FID:{self.file_id}] Skipping element extraction for Slide {slide_index_1based} due to pptx slide object unavailability.")


                    slide_processed_successfully = True

                except Exception as slide_e:
                     logger.error(f"[Converter FID:{self.file_id}] Failed to process slide {slide_index_1based}: {slide_e}", exc_info=True)
                     encountered_error = True # Mark that an error occurred during this file's conversion
                     # Optional: break here to stop processing this file? Or continue to next slide?
                     # Let's continue for now to try and process other slides.

                # --- 4e. Emit Progress ---
                processed_slides += 1
                self.signals.progress.emit(self.file_id, slide_index_1based, total_slides)


            # --- 5. Update Final File Status ---
            if self.is_cancelled:
                 final_status = 'Pending' # Revert to Pending if cancelled? Or keep In Progress? Let's use Pending.
                 error_message = "Cancelled"
            elif encountered_error:
                 final_status = 'Failed'
                 error_message = "Failed to process one or more slides. Check logs."
            else:
                 final_status = 'Completed'
                 error_message = "" # No error

            logger.info(f"[Converter FID:{self.file_id}] Setting final status to: {final_status}")
            thread_local_db.update_file_conversion_status(self.file_id, final_status)

            # --- 6. Emit Final Signal ---
            if final_status == 'Completed':
                 self.signals.finished.emit(self.file_id)
            else:
                 # Pass the aggregated/final error message
                 self.signals.error.emit(self.file_id, error_message)


        except Exception as e:
            # Catch major errors (COM open, pptx open, initial setup)
            logger.error(f"[Converter FID:{self.file_id}] Conversion worker failed critically: {e}", exc_info=True)
            error_message = f"Critical conversion error: {e}"
            try:
                # Try to mark the file as failed in DB if a critical error occurred
                thread_local_db.update_file_conversion_status(self.file_id, 'Failed')
            except Exception as db_e:
                logger.error(f"[Converter FID:{self.file_id}] Failed to update DB status to Failed after critical error: {db_e}")
            self.signals.error.emit(self.file_id, error_message)

        finally:
            # --- 7. Cleanup Resources ---
            # Close thread-local database connection
            if 'thread_local_db' in locals():
                try:
                    thread_local_db.close()
                    logger.debug(f"[Converter FID:{self.file_id}] Closed thread-local database connection")
                except Exception as db_close_e:
                    logger.error(f"[Converter FID:{self.file_id}] Error closing thread-local database: {db_close_e}", exc_info=True)
            
            # Close COM objects
            if 'presentation_com' in locals() and presentation_com is not None:
                try:
                    presentation_com.Close()
                    logger.debug(f"[Converter FID:{self.file_id}] Closed PowerPoint presentation")
                except Exception as e:
                    logger.error(f"[Converter FID:{self.file_id}] Error closing presentation: {e}", exc_info=True)
            
            # Quit PowerPoint application to prevent process leaks
            if 'ppt_app' in locals() and ppt_app is not None:
                try:
                    ppt_app.Quit()
                    logger.debug(f"[Converter FID:{self.file_id}] Quit PowerPoint application")
                except Exception as e:
                    logger.error(f"[Converter FID:{self.file_id}] Error quitting PowerPoint: {e}", exc_info=True)
            
            # Important: Release COM objects to allow PowerPoint to exit eventually
            presentation_com = None
            ppt_app = None
            try:
                pythoncom.CoUninitialize()
                logger.debug(f"[Converter FID:{self.file_id}] CoUninitialize called.")
            except Exception as e:
                logger.warning(f"[Converter FID:{self.file_id}] Error in CoUninitialize: {e}")

        logger.info(f"[Converter FID:{self.file_id}] Worker finished execution.")

    def cancel(self):
        """Requests cancellation of the conversion process."""
        logger.info(f"[Converter FID:{self.file_id}] Cancellation requested.")
        self.is_cancelled = True