# src/slideman/services/slide_converter.py

import logging
import time
import traceback
from pathlib import Path
from typing import Optional

# --- Qt Imports ---
from PySide6.QtCore import QObject, QRunnable, Signal, Slot, QSize, Qt
from PySide6.QtGui import QPixmap

# --- External Libraries ---
try:
    import pythoncom
    import win32com.client
    import pywintypes
    HAS_COM = True
except ImportError:
    HAS_COM = False
    logging.error("pywin32 library not found. PowerPoint COM automation will not work.")
    pywintypes = type('pywintypes', (object,), {'com_error': Exception})()

try:
    from pptx import Presentation
    from pptx.util import Inches, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logging.error("python-pptx library not found. Shape extraction will not work.")
    MSO_SHAPE_TYPE = type('MSO_SHAPE_TYPE', (object,), {
        'PLACEHOLDER': 0, 'PICTURE': 1, 'CHART': 2, 'TABLE': 3, 
        'TEXT_BOX': 4, 'GROUP': 5, 'AUTO_SHAPE': 6, 'FREEFORM': 7
    })()

# --- Local Imports ---
from .database_worker import DatabaseWorker
from .exceptions import (
    DatabaseError, PowerPointError, COMInitializationError,
    PresentationAccessError, SlideExportError, ValidationError
)

# --- Constants ---
THUMBNAIL_HEIGHT = 200  # Target height in pixels for thumbnails

logger = logging.getLogger(__name__)


class SlideConverterSignals(QObject):
    """Defines signals emitted by the SlideConverter worker."""
    progress = Signal(int, int, int)  # file_id, current_slide_index (1-based), total_slides
    finished = Signal(int)            # file_id
    error = Signal(int, str)         # file_id, error_message


class SlideConverter(QRunnable):
    """
    Thread-safe worker to convert PowerPoint slides to images,
    extract shape data, and store results in the database.
    """
    
    def __init__(self, file_id: int, file_path: Path, db_path: Path, signals: SlideConverterSignals):
        """
        Initialize the slide converter.
        
        Args:
            file_id: Database ID of the file to convert.
            file_path: Path to the PowerPoint file.
            db_path: Path to the database for worker connection.
            signals: Signals object for progress updates.
            
        Raises:
            ImportError: If required libraries are not available.
        """
        super().__init__()
        if not HAS_COM or not HAS_PPTX:
            raise ImportError("Missing required libraries (pywin32 or python-pptx) for SlideConverter.")

        self.signals = signals
        self.file_id = file_id
        self.file_path = file_path
        self.db_path = db_path
        self.is_cancelled = False

        # Determine paths for converted data
        self.project_root = file_path.parent.parent
        self.converted_data_dir = self.project_root / "converted_data" / str(self.file_id)

    def map_shape_type(self, pptx_shape_type) -> str:
        """Maps python-pptx MSO_SHAPE_TYPE to simple strings."""
        if pptx_shape_type in (MSO_SHAPE_TYPE.PLACEHOLDER, MSO_SHAPE_TYPE.AUTO_SHAPE,
                               MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.GROUP):
            return "SHAPE"
        elif pptx_shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "PICTURE"
        elif pptx_shape_type == MSO_SHAPE_TYPE.CHART:
            return "CHART"
        elif pptx_shape_type == MSO_SHAPE_TYPE.TABLE:
            return "TABLE"
        elif pptx_shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            return "TEXT"
        else:
            return "UNKNOWN"

    @Slot()
    def run(self):
        """Main execution logic running in the background thread."""
        logger.info(f"[Converter FID:{self.file_id}] Starting conversion for: {self.file_path}")

        # Initialize resources
        db_worker = None
        ppt_app = None
        presentation_com = None
        com_initialized = False
        
        try:
            # Initialize COM for this thread
            try:
                pythoncom.CoInitialize()
                com_initialized = True
                logger.debug(f"[Converter FID:{self.file_id}] COM Initialized.")
            except Exception as e:
                raise COMInitializationError(f"Failed to initialize COM: {e}") from e

            # Create database worker
            db_worker = DatabaseWorker(self.db_path)
            
            # Update conversion status to In Progress
            try:
                db_worker.update_file_conversion_status(self.file_id, "In Progress")
            except DatabaseError as e:
                logger.warning(f"[Converter FID:{self.file_id}] Failed to update initial status: {e}")

            # Create output directories
            self.converted_data_dir.mkdir(parents=True, exist_ok=True)
            logger.debug(f"[Converter FID:{self.file_id}] Output directory ready: {self.converted_data_dir}")

            # Process the presentation
            self._process_presentation(db_worker)

        except COMInitializationError as e:
            logger.error(f"[Converter FID:{self.file_id}] COM initialization failed: {e}")
            self.signals.error.emit(self.file_id, str(e))
            self._update_status_safe(db_worker, "Failed")
            
        except PowerPointError as e:
            logger.error(f"[Converter FID:{self.file_id}] PowerPoint error: {e}")
            self.signals.error.emit(self.file_id, str(e))
            self._update_status_safe(db_worker, "Failed")
            
        except Exception as e:
            logger.error(f"[Converter FID:{self.file_id}] Unexpected error: {e}", exc_info=True)
            self.signals.error.emit(self.file_id, f"Conversion failed: {str(e)}")
            self._update_status_safe(db_worker, "Failed")
            
        finally:
            # Cleanup resources
            self._cleanup(db_worker, presentation_com, ppt_app, com_initialized)
            
        logger.info(f"[Converter FID:{self.file_id}] Worker finished.")

    def _process_presentation(self, db_worker: DatabaseWorker):
        """
        Process the PowerPoint presentation.
        
        Args:
            db_worker: Database worker instance.
            
        Raises:
            PowerPointError: If presentation cannot be processed.
        """
        ppt_app = None
        presentation_com = None
        presentation_pptx = None
        
        try:
            # Open presentation with COM
            logger.debug(f"[Converter FID:{self.file_id}] Opening presentation via COM...")
            try:
                ppt_app = win32com.client.Dispatch("PowerPoint.Application")
                presentation_com = ppt_app.Presentations.Open(
                    str(self.file_path), ReadOnly=True, WithWindow=False
                )
                total_slides = len(presentation_com.Slides)
                logger.info(f"[Converter FID:{self.file_id}] Opened via COM. Found {total_slides} slides.")
            except pywintypes.com_error as e:
                raise PresentationAccessError(
                    f"Cannot open PowerPoint file: {self.file_path.name}"
                ) from e

            # Update slide count in database
            try:
                db_worker.update_file_slide_count(self.file_id, total_slides)
            except DatabaseError as e:
                logger.warning(f"[Converter FID:{self.file_id}] Failed to update slide count: {e}")

            # Open with python-pptx for shape extraction
            logger.debug(f"[Converter FID:{self.file_id}] Opening presentation via python-pptx...")
            try:
                presentation_pptx = Presentation(self.file_path)
                if total_slides != len(presentation_pptx.slides):
                    logger.warning(
                        f"[Converter FID:{self.file_id}] Slide count mismatch! "
                        f"COM: {total_slides}, python-pptx: {len(presentation_pptx.slides)}"
                    )
            except Exception as e:
                logger.error(f"[Converter FID:{self.file_id}] python-pptx error: {e}")
                presentation_pptx = None  # Continue without shape extraction

            # Handle empty presentations
            if total_slides == 0:
                logger.warning(f"[Converter FID:{self.file_id}] Presentation has no slides.")
                db_worker.update_file_conversion_status(self.file_id, "Completed")
                self.signals.finished.emit(self.file_id)
                return

            # Process each slide
            slides_processed = 0
            errors_encountered = []
            
            for i in range(total_slides):
                if self.is_cancelled:
                    logger.info(f"[Converter FID:{self.file_id}] Conversion cancelled.")
                    raise ValidationError("Conversion cancelled by user")

                slide_index = i + 1
                try:
                    self._process_slide(
                        db_worker, presentation_com, presentation_pptx, 
                        slide_index, total_slides
                    )
                    slides_processed += 1
                    
                except Exception as e:
                    error_msg = f"Failed to process slide {slide_index}: {str(e)}"
                    logger.error(f"[Converter FID:{self.file_id}] {error_msg}")
                    errors_encountered.append(error_msg)
                    # Continue processing other slides
                
                # Emit progress
                self.signals.progress.emit(self.file_id, slide_index, total_slides)

            # Determine final status
            if errors_encountered:
                if slides_processed == 0:
                    raise SlideExportError("Failed to process any slides")
                else:
                    # Partial success
                    logger.warning(
                        f"[Converter FID:{self.file_id}] Processed {slides_processed}/{total_slides} slides. "
                        f"Errors: {len(errors_encountered)}"
                    )
                    db_worker.update_file_conversion_status(self.file_id, "Completed")
                    self.signals.finished.emit(self.file_id)
            else:
                # Complete success
                logger.info(f"[Converter FID:{self.file_id}] Successfully processed all {total_slides} slides.")
                db_worker.update_file_conversion_status(self.file_id, "Completed")
                self.signals.finished.emit(self.file_id)
                
        finally:
            # Cleanup COM objects
            if presentation_com:
                try:
                    presentation_com.Close()
                except Exception as e:
                    logger.error(f"[Converter FID:{self.file_id}] Error closing presentation: {e}")
            if ppt_app:
                try:
                    ppt_app.Quit()
                except Exception as e:
                    logger.error(f"[Converter FID:{self.file_id}] Error quitting PowerPoint: {e}")

    def _process_slide(self, db_worker: DatabaseWorker, presentation_com, 
                      presentation_pptx, slide_index: int, total_slides: int):
        """
        Process a single slide.
        
        Args:
            db_worker: Database worker instance.
            presentation_com: COM presentation object.
            presentation_pptx: python-pptx presentation object.
            slide_index: 1-based slide index.
            total_slides: Total number of slides.
            
        Raises:
            SlideExportError: If slide processing fails.
        """
        logger.debug(f"[Converter FID:{self.file_id}] Processing slide {slide_index}/{total_slides}")
        
        try:
            # Get slide objects
            slide_com = presentation_com.Slides(slide_index)
            slide_pptx = None
            if presentation_pptx and (slide_index - 1) < len(presentation_pptx.slides):
                slide_pptx = presentation_pptx.slides[slide_index - 1]

            # Export full-resolution image
            img_filename = f"image_{slide_index}.png"
            img_full_path = self.converted_data_dir / img_filename
            
            try:
                slide_com.Export(str(img_full_path), "PNG")
                image_rel_path = img_full_path.relative_to(self.project_root).as_posix()
            except Exception as e:
                raise SlideExportError(f"Failed to export slide image: {e}") from e

            # Create thumbnail
            thumb_filename = f"thumb_{slide_index}.png"
            thumb_full_path = self.converted_data_dir / thumb_filename
            thumb_rel_path = self._create_thumbnail(img_full_path, thumb_full_path)

            # Add slide to database
            slide_id = db_worker.add_slide_with_paths(
                self.file_id, slide_index, thumb_rel_path, image_rel_path
            )

            # Extract and store elements if available
            if slide_pptx:
                self._extract_elements(db_worker, slide_id, slide_pptx)
                
        except Exception as e:
            raise SlideExportError(f"Slide {slide_index} processing failed: {e}") from e

    def _create_thumbnail(self, source_path: Path, thumb_path: Path) -> Optional[str]:
        """
        Create a thumbnail from the source image.
        
        Args:
            source_path: Path to source image.
            thumb_path: Path for thumbnail output.
            
        Returns:
            Relative path to thumbnail or None if failed.
        """
        try:
            pixmap = QPixmap(str(source_path))
            if pixmap.isNull():
                logger.error(f"[Converter FID:{self.file_id}] Failed to load image: {source_path}")
                return None

            # Scale to desired height, keeping aspect ratio
            thumb_pixmap = pixmap.scaledToHeight(
                THUMBNAIL_HEIGHT, Qt.TransformationMode.SmoothTransformation
            )
            
            if not thumb_pixmap.save(str(thumb_path), "PNG"):
                logger.error(f"[Converter FID:{self.file_id}] Failed to save thumbnail: {thumb_path}")
                return None
                
            return thumb_path.relative_to(self.project_root).as_posix()
            
        except Exception as e:
            logger.error(f"[Converter FID:{self.file_id}] Thumbnail creation failed: {e}")
            return None

    def _extract_elements(self, db_worker: DatabaseWorker, slide_id: int, slide_pptx):
        """
        Extract elements from a slide using python-pptx.
        
        Args:
            db_worker: Database worker instance.
            slide_id: Database ID of the slide.
            slide_pptx: python-pptx slide object.
        """
        element_count = 0
        
        for shape in slide_pptx.shapes:
            try:
                # Check if shape has required attributes
                if not all(hasattr(shape, attr) for attr in ['shape_type', 'left', 'top', 'width', 'height']):
                    continue
                    
                if any(getattr(shape, attr) is None for attr in ['left', 'top', 'width', 'height']):
                    continue

                # Extract shape data
                element_type = self.map_shape_type(shape.shape_type)
                
                # Extract text content if available
                content = ""
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    content = shape.text.strip()

                # Add element to database
                element_id = db_worker.add_element(slide_id, element_type, content)
                if element_id:
                    element_count += 1
                    
            except Exception as e:
                logger.warning(
                    f"[Converter FID:{self.file_id}] Failed to extract element: {e}"
                )
                
        logger.debug(f"[Converter FID:{self.file_id}] Extracted {element_count} elements from slide {slide_id}")

    def _update_status_safe(self, db_worker: Optional[DatabaseWorker], status: str):
        """Safely update file status, handling errors."""
        if db_worker:
            try:
                db_worker.update_file_conversion_status(self.file_id, status)
            except Exception as e:
                logger.error(f"[Converter FID:{self.file_id}] Failed to update status to {status}: {e}")

    def _cleanup(self, db_worker: Optional[DatabaseWorker], 
                presentation_com, ppt_app, com_initialized: bool):
        """Clean up all resources."""
        # Close database connection
        if db_worker:
            try:
                db_worker.close()
                logger.debug(f"[Converter FID:{self.file_id}] Closed database connection")
            except Exception as e:
                logger.error(f"[Converter FID:{self.file_id}] Error closing database: {e}")

        # Cleanup COM
        if com_initialized:
            try:
                pythoncom.CoUninitialize()
                logger.debug(f"[Converter FID:{self.file_id}] COM uninitialized")
            except Exception as e:
                logger.warning(f"[Converter FID:{self.file_id}] Error uninitializing COM: {e}")

    def cancel(self):
        """Request cancellation of the conversion process."""
        logger.info(f"[Converter FID:{self.file_id}] Cancellation requested.")
        self.is_cancelled = True